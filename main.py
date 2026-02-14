import pandas as pd
import feedparser
import yfinance as yf
import os
import io
import requests
from datetime import datetime
from pptx import Presentation
from pptx.util import Pt

def get_stock_universe():
    url = "https://www.niftyindices.com/IndexConstituent/ind_nifty500list.csv"
    headers = {'User-Agent': 'Mozilla/5.0'}
    try:
        response = requests.get(url, headers=headers, timeout=15)
        df = pd.read_csv(io.StringIO(response.text))
        return df
    except Exception as e:
        print(f"⚠️ Error fetching universe: {e}")
        return pd.DataFrame()

def get_detailed_analysis(sector, change):
    # Sectoral Intelligence Dictionary for 2026
    logic = {
        "Financial Services": {
            "time": "6-12 Months", "upside": "12-15%",
            "reason": "Strong credit growth and cleaner balance sheets in the NBFC space."
        },
        "Healthcare": {
            "time": "12-24 Months", "upside": "20%",
            "reason": "Shift toward digital healthcare and specialized pharma therapies."
        },
        "Energy": {
            "time": "3-5 Years", "upside": "25%+",
            "reason": "Massive capital rotation into Green Hydrogen and Renewable infrastructure."
        },
        "Capital Goods": {
            "time": "18-36 Months", "upside": "18%",
            "reason": "Peak manufacturing cycles driven by matured PLI incentive schemes."
        }
    }
    info = logic.get(sector, {"time": "12 Months", "upside": "10-12%", "reason": "Stable growth tracking Nifty index earnings."})
    sentiment = "Bullish" if change > 0 else "Consolidating"
    return (f"Sentiment: {sentiment}\nTarget Growth: {info['upside']} | Timeframe: {info['time']}\nContext: {info['reason']}")

def create_ppt(data_list, filename):
    os.makedirs("reports", exist_ok=True)
    filepath = os.path.join("reports", filename)
    prs = Presentation()
    
    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Strategic Sectoral Analysis"
    slide.placeholders[1].text = f"Growth & Timeframe Forecasts\nGenerated: {datetime.now().strftime('%d %b %Y')}"

    for d in data_list:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"{d['name']} ({d['sector']})"
        body = slide.placeholders[1].text_frame
        
        p1 = body.paragraphs[0]
        p1.text = f"Current Performance: {d['change']}%"
        p1.font.bold = True
        
        p2 = body.add_paragraph()
        p2.text = f"\nMarket News:\n{d['headline']}"
        p2.font.size = Pt(12)

        p3 = body.add_paragraph()
        p3.text = f"\nDeep-Dive Analysis:\n{d['analysis']}"
        p3.font.size = Pt(14)
        p3.font.bold = True

    prs.save(filepath)
    print(f"✅ Saved to {filepath}")

def run_analysis():
    # FIX: Initialize variable to avoid NameError if no news matches
    final_data = [] 
    df = get_stock_universe()
    if df.empty: return

    stock_map = {row['Symbol']: row['Company Name'] for _, row in df.iterrows()}
    sector_map = {row['Symbol']: row.get('Industry', 'General') for _, row in df.iterrows()}

    feed = feedparser.parse("https://news.google.com/rss/search?q=Indian+stock+market+news&hl=en-IN&gl=IN&ceid=IN:en")
    seen = set()

    for entry in feed.entries[:100]:
        if len(final_data) >= 15: break
        headline = entry.title
        for sym, name in stock_map.items():
            if sym in headline or name.split()[0] in headline:
                ticker = f"{sym}.NS"
                if ticker not in seen:
                    try:
                        s = yf.Ticker(ticker).history(period="2d")
                        if len(s) < 2: continue
                        change = round(((s['Close'].iloc[-1]-s['Close'].iloc[-2])/s['Close'].iloc[-2])*100, 2)
                        final_data.append({
                            "name": name, "sector": sector_map[sym], "headline": headline,
                            "change": change, "analysis": get_detailed_analysis(sector_map[sym], change)
                        })
                        seen.add(ticker)
                        break
                    except: continue

    if final_data:
        date_str = datetime.now().strftime('%d_%b_%Y')
        filename = f"Market_DeepDive_{date_str}.pptx"
        create_ppt(final_data, filename)
    else:
        print("No stock news found for deep dive today.")

if __name__ == "__main__":
    run_analysis()
